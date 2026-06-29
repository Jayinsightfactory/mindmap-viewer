# -*- coding: utf-8 -*-
"""
엑셀 → PPT 품목 카탈로그 생성기 (Orbit AI / nenova)  v2

핵심
  - 슬라이드를 여러 개 생성, 슬라이드마다 가로×세로(그리드)를 개별 지정
  - 슬라이드마다 채울 품목을 개별 선택
  - 작업 내용(품목·이미지·슬라이드 구성) 자동 저장 → 재실행 시 복원
  - 품목 라이브러리: 큰 썸네일 + 4배 큰 체크박스, 드래그 복수선택
필요 패키지:  pip install openpyxl python-pptx pillow
exe 빌드   :  build.bat 실행 (PyInstaller --noconsole)
"""

import io
import os
import sys
import json
import zipfile
import subprocess

import openpyxl
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor

# ---------------- 슬라이드 크기 (템플릿 기준 27.5 x 19.1 cm) ----------------
SLIDE_W = 9906000
SLIDE_H = 6858000

KW_NAME   = ("품목", "품명", "name", "item")
KW_ORIGIN = ("원산지", "origin")
KW_SEASON = ("공급", "기간", "시즌", "season")
KW_SIZE   = ("사이즈", "규격", "size")
KW_PRICE  = ("도착원가", "도착 원가", "도착", "원가", "단가")

# 슬라이드에 넣을 수 있는 항목 (라벨 / 기본 순서)
FIELD_LABELS = {
    "name":   "품목명",
    "size":   "사이즈",
    "season": "공급가능기간",
    "origin": "원산지",
    "price":  "도착원가(부가세포함)",
}
DEFAULT_FIELDS = ["name", "size", "season", "origin", "price"]
ARRIVAL_DIVISOR = 0.85          # 도착원가 입력 시 ÷0.85

# 저장 위치 (작업 자동 저장)
APP_DIR   = os.path.join(os.path.expanduser("~"), ".orbit_ppt")
CACHE_DIR = os.path.join(APP_DIR, "cache")
PROJECT   = os.path.join(APP_DIR, "project.json")


def _s(v):
    return (str(v) if v is not None else "").strip()


import re as _re


def _to_number(v):
    """숫자/숫자형 문자열(콤마·통화 포함)을 float로. 실패 시 None."""
    if v is None:
        return None
    if isinstance(v, (int, float)):
        return float(v)
    s = str(v).replace(",", "")
    m = _re.search(r"-?\d+(?:\.\d+)?", s)
    return float(m.group()) if m else None


# ==================================================================
# 그리드 슬롯 계산 (가로 cols × 세로 rows → 셀마다 이미지+텍스트 위치)
# ==================================================================
def grid_slots(cols, rows, top_pad=320000, slide_w=SLIDE_W, slide_h=SLIDE_H):
    mx, mb, gut = 500000, 300000, 160000
    my = top_pad
    usable_w = slide_w - 2 * mx
    usable_h = slide_h - my - mb
    cw = (usable_w - (cols - 1) * gut) / cols
    ch = (usable_h - (rows - 1) * gut) / rows
    slots = []
    for r in range(rows):
        for c in range(cols):
            cl = mx + c * (cw + gut)
            ct = my + r * (ch + gut)
            txt_h = min(ch * 0.30, 1000000)
            gap = 40000
            img_box = max(200000, min(cw, ch - txt_h - gap))
            img_l = cl + (cw - img_box) / 2
            slots.append((
                int(img_l), int(ct), int(img_box),
                int(cl), int(ct + img_box + gap),
                int(cw), int(ch - img_box - gap)))
    return slots


# ==================================================================
# 데이터 모델 + 엑셀 추출
# ==================================================================
class Item:
    __slots__ = ("name", "origin", "season", "size", "price",
                 "img_bytes", "sheet", "row")

    def __init__(self, name, origin, season, size, price,
                 img_bytes, sheet, row):
        self.name = _s(name)
        self.origin = _s(origin)
        self.season = _s(season)
        self.size = _s(size)
        self.price = price             # 도착원가 원본값(숫자) 또는 None
        self.img_bytes = img_bytes
        self.sheet = sheet
        self.row = row

    @property
    def key(self):
        return (self.sheet, self.row)


# ----- 항목별 텍스트 생성 (체크/순서/도착원가 계산/공급기간 치환 반영) -----
def field_text(item, key, season_map=None):
    if key == "name":
        return item.name
    if key == "size":
        return item.size
    if key == "origin":
        return item.origin
    if key == "season":
        v = item.season
        if season_map:
            v = season_map.get(v.strip(), v)   # 사용자 정의 치환
        return v
    if key == "price":
        p = _to_number(item.price)
        if p is None or p <= 0:
            return ""
        return f"{round(p / ARRIVAL_DIVISOR):,}원"
    return ""


def build_lines(item, fields=None, season_map=None):
    fields = fields or DEFAULT_FIELDS
    out = []
    for k in fields:
        v = field_text(item, k, season_map)
        if v:
            out.append(f"{FIELD_LABELS.get(k, k)} : {v}")
    return out


def _find_col(ws, header_row, keywords):
    kws = [k.lower() for k in keywords]
    for c in range(1, ws.max_column + 1):
        v = ws.cell(row=header_row, column=c).value
        if v and any(k in str(v).lower() for k in kws):
            return c
    return None


def _find_col_rows(ws, rows, keywords):
    kws = [k.lower() for k in keywords]
    for hr in rows:
        for c in range(1, ws.max_column + 1):
            v = ws.cell(row=hr, column=c).value
            if v and any(k in str(v).lower() for k in kws):
                return c
    return None


def _detect_header_row(ws):
    kws = [k.lower() for k in KW_NAME]
    for r in range(1, min(6, ws.max_row + 1)):
        for c in range(1, ws.max_column + 1):
            v = ws.cell(row=r, column=c).value
            if v and any(k in str(v).lower() for k in kws):
                return r
    return 1


def _merge_map(ws):
    """병합셀 → 모든 칸이 좌상단 값을 갖도록 매핑 (사이즈/공급기간 병합 대응)."""
    m = {}
    for rng in ws.merged_cells.ranges:
        tl = ws.cell(row=rng.min_row, column=rng.min_col).value
        if tl is None:
            continue
        for r in range(rng.min_row, rng.max_row + 1):
            for c in range(rng.min_col, rng.max_col + 1):
                m[(r, c)] = tl
    return m


def _cell(ws, mm, r, c):
    if not c:
        return None
    v = ws.cell(row=r, column=c).value
    if v is None:
        v = mm.get((r, c))
    return v


# 마지막 load_items 진단 정보 (이미지 인식 개수 등)
LOAD_DIAG = {"items": 0, "images_total": 0, "images_mapped": 0, "images_used": 0}


def _img_anchor_row(im):
    a = getattr(im, "anchor", None)
    frm = getattr(a, "_from", None)
    if frm is not None and hasattr(frm, "row"):
        try:
            return int(frm.row) + 1
        except Exception:
            return None
    return None


def _safe_img_data(im):
    try:
        return im._data()
    except Exception:
        try:
            ref = getattr(im, "ref", None)
            if hasattr(ref, "getvalue"):
                return ref.getvalue()
        except Exception:
            pass
    return None


def _extract_images(ws, path):
    """이미지를 행 기준으로 추출. 앵커 매핑 실패 시 순서 배정 / zip 미디어 폴백."""
    mapped, ordered = {}, []
    imgs = list(getattr(ws, "_images", []))
    for im in imgs:
        data = _safe_img_data(im)
        if not data:
            continue
        ordered.append(data)
        row = _img_anchor_row(im)
        if row is not None:
            mapped.setdefault(row, data)
    total = len(imgs)
    if not ordered:                       # 워크시트에서 이미지를 전혀 못 읽음 → zip 직접
        try:
            with zipfile.ZipFile(path) as z:
                for n in sorted(z.namelist()):
                    low = n.lower()
                    if n.startswith("xl/media/") and low.split(".")[-1] in (
                            "png", "jpg", "jpeg", "gif", "bmp", "emf", "wmf"):
                        ordered.append(z.read(n))
        except Exception:
            pass
        total = max(total, len(ordered))
    return mapped, ordered, total


def list_sheets(path):
    wb = openpyxl.load_workbook(path, read_only=False)
    names = wb.sheetnames
    wb.close()
    return names


def load_items(path, sheet_name):
    # data_only=True: 수식 셀의 캐시된 계산값을 읽음(도착원가 등)
    wb = openpyxl.load_workbook(path, data_only=True)
    ws = wb[sheet_name]
    header_row = _detect_header_row(ws)
    hrows = list(range(header_row, min(header_row + 3, ws.max_row + 1)))
    name_col   = _find_col(ws, header_row, KW_NAME)   or 1
    origin_col = _find_col(ws, header_row, KW_ORIGIN)
    season_col = _find_col(ws, header_row, KW_SEASON)
    size_col   = _find_col(ws, header_row, KW_SIZE)
    price_col  = _find_col_rows(ws, hrows, KW_PRICE)   # 도착원가(2행 헤더)
    mm = _merge_map(ws)

    # 헤더 블록 끝(병합 헤더 B1:B2 / FOB·도착원가 서브헤더) 다음부터 데이터
    block_end = header_row
    for rng in ws.merged_cells.ranges:
        if rng.min_row == header_row:
            block_end = max(block_end, rng.max_row)
    sub_col = _find_col_rows(ws, range(header_row + 1, header_row + 3),
                             KW_PRICE + ("fob",))
    if sub_col:
        for hr in range(header_row + 1, header_row + 3):
            v = ws.cell(row=hr, column=sub_col).value
            if v and any(k in str(v).lower() for k in
                         [x.lower() for x in KW_PRICE + ("fob",)]):
                block_end = max(block_end, hr)

    mapped, ordered, total_imgs = _extract_images(ws, path)
    # 앵커 매핑이 전혀 안 되면 이름 있는 행에 순서대로 이미지 배정
    use_seq = (len(mapped) == 0 and len(ordered) > 0)
    seq = iter(ordered)

    items, used = [], 0
    for r in range(block_end + 1, ws.max_row + 1):
        name = _cell(ws, mm, r, name_col)
        img = mapped.get(r)
        if img is None and use_seq and name:
            img = next(seq, None)
        if not name and not img:
            continue
        if not name:
            name = f"(이름없음 {r}행)"
        if img:
            used += 1
        items.append(Item(
            name,
            _cell(ws, mm, r, origin_col),
            _cell(ws, mm, r, season_col),
            _cell(ws, mm, r, size_col),
            _cell(ws, mm, r, price_col),
            img, sheet_name, r))
    wb.close()
    LOAD_DIAG.update(items=len(items), images_total=total_imgs,
                     images_mapped=len(mapped), images_used=used)
    return items


# ==================================================================
# PPT 생성 (슬라이드마다 그리드/품목 다름)
# ==================================================================
def _fit(img_bytes, box):
    try:
        im = Image.open(io.BytesIO(img_bytes))
        iw, ih = im.size
    except Exception:
        return box, box, 0, 0
    s = min(box / iw, box / ih)
    w, h = int(iw * s), int(ih * s)
    return w, h, (box - w) // 2, (box - h) // 2


def _fit_font_pt(n_lines, box_h_emu, base_pt):
    """줄 수 × 박스 높이에 맞춰 글자 크기(pt) 자동 축소 — 모든 항목이 칸에 들어가게."""
    n = max(1, n_lines)
    max_pt = (box_h_emu / n) / 12700 / 1.3      # 1pt=12700EMU, 줄간격 여유 1.3
    return max(7, int(min(base_pt, max_pt)))


def build_pptx(slides, store, out_path, fields=None, name_pt=14,
               titles=None, season_map=None,
               slide_w=SLIDE_W, slide_h=SLIDE_H):
    """
    slides : [{"cols":3,"rows":2,"items":[key, key, ...]}, ...]
    store  : {key(tuple): Item}
    titles : 슬라이드별 제목 리스트 또는 None
    """
    prs = Presentation()
    prs.slide_width = slide_w
    prs.slide_height = slide_h
    blank = prs.slide_layouts[6]
    from pptx.enum.shapes import MSO_SHAPE

    for si, sd in enumerate(slides):
        cols, rows = int(sd["cols"]), int(sd["rows"])
        slide = prs.slides.add_slide(blank)

        title = titles[si] if (titles and si < len(titles) and titles[si]) else None
        top_pad = 320000
        if title:
            top_pad = 880000
            tb = slide.shapes.add_textbox(
                Emu(500000), Emu(180000), Emu(slide_w - 1000000), Emu(560000))
            p = tb.text_frame.paragraphs[0]
            run = p.add_run()
            run.text = str(title)
            run.font.size = Pt(20)
            run.font.bold = True

        slots = grid_slots(cols, rows, top_pad, slide_w, slide_h)
        keys = sd["items"][:cols * rows]
        for idx, key in enumerate(keys):
            it = store.get(tuple(key))
            if not it:
                continue
            il, it_top, ibox, tl, tt, tw, th = slots[idx]
            if it.img_bytes:
                w, h, ox, oy = _fit(it.img_bytes, ibox)
                try:
                    slide.shapes.add_picture(
                        io.BytesIO(it.img_bytes),
                        Emu(il + ox), Emu(it_top + oy), Emu(w), Emu(h))
                except Exception:
                    pass
            else:
                ph = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE, Emu(il), Emu(it_top),
                    Emu(ibox), Emu(ibox))
                ph.fill.solid()
                ph.fill.fore_color.rgb = RGBColor(0xEE, 0xEE, 0xEE)
                ph.line.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)

            box_h = max(th, 300000)
            box = slide.shapes.add_textbox(Emu(tl), Emu(tt), Emu(tw), Emu(box_h))
            tf = box.text_frame
            tf.word_wrap = True
            lines = build_lines(it, fields, season_map)
            fpt = _fit_font_pt(len(lines), box_h, name_pt)
            for li, line in enumerate(lines):
                p = tf.paragraphs[0] if li == 0 else tf.add_paragraph()
                run = p.add_run()
                run.text = line
                run.font.size = Pt(fpt)
                run.font.bold = (li == 0)

    prs.save(out_path)
    return len(slides)


# ==================================================================
# 슬라이드 미리보기 렌더 (PIL — tkinter 없이 사용 가능)
# ==================================================================
_FONT_CANDIDATES = [
    "C:/Windows/Fonts/malgun.ttf",
    "C:/Windows/Fonts/malgunsl.ttf",
    "/usr/share/fonts/truetype/nanum/NanumGothic.ttf",
    "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
]


def _load_font(size):
    for p in _FONT_CANDIDATES:
        if os.path.exists(p):
            try:
                return ImageFont.truetype(p, size)
            except Exception:
                pass
    return ImageFont.load_default()


def render_slide_image(sd, store, title=None, fields=None,
                       season_map=None, width=940,
                       slide_w=SLIDE_W, slide_h=SLIDE_H):
    """슬라이드 1장을 PIL 이미지로 렌더 (미리보기용)."""
    cols, rows = int(sd["cols"]), int(sd["rows"])
    sc = slide_w / width
    W, H = width, int(slide_h / sc)
    cv = Image.new("RGB", (W, H), "white")
    d = ImageDraw.Draw(cv)
    d.rectangle([0, 0, W - 1, H - 1], outline="#cccccc")

    def px(v):
        return int(v / sc)

    top_pad = 320000
    if title:
        top_pad = 880000
        d.text((px(500000), px(220000)), str(title),
               fill="black", font=_load_font(max(14, px(360000))))

    slots = grid_slots(cols, rows, top_pad, slide_w, slide_h)
    keys = sd["items"][:cols * rows]
    for i, key in enumerate(keys):
        it = store.get(tuple(key))
        il, it_top, ibox, tl, tt, tw, th = slots[i]
        d.rectangle([px(il), px(it_top), px(il + ibox), px(it_top + ibox)],
                    outline="#dddddd")
        if it and it.img_bytes:
            w, h, ox, oy = _fit(it.img_bytes, ibox)
            try:
                im = Image.open(io.BytesIO(it.img_bytes)).convert("RGB")
                im = im.resize((max(1, px(w)), max(1, px(h))))
                cv.paste(im, (px(il + ox), px(it_top + oy)))
            except Exception:
                pass
        if it:
            lines = build_lines(it, fields, season_map)
            avail = max(1, px(th))
            fs = max(8, min(int(px(ibox) * 0.11),
                            int(avail / max(1, len(lines)) / 1.25)))
            font = _load_font(fs)
            maxw = px(tw) - 4
            y = px(tt) + 1
            for line in lines:
                for seg in _wrap(d, line, font, maxw):
                    d.text((px(tl) + 2, y), seg, fill="black", font=font)
                    y += fs + 2
    return cv


def _wrap(draw, text, font, maxw):
    """maxw(px) 안에 들어가도록 단어/글자 단위 줄바꿈."""
    def w(s):
        return draw.textlength(s, font=font)
    if w(text) <= maxw:
        return [text]
    out, cur = [], ""
    for ch in text:
        if w(cur + ch) > maxw and cur:
            out.append(cur)
            cur = ch
        else:
            cur += ch
    if cur:
        out.append(cur)
    return out


def open_file(path):
    """OS 기본 프로그램으로 파일 열기 (Windows에서 PPT 자동 실행)."""
    try:
        if os.name == "nt":
            os.startfile(path)            # Windows
        elif sys.platform == "darwin":
            subprocess.Popen(["open", path])
        else:
            subprocess.Popen(["xdg-open", path])
        return True
    except Exception:
        return False


# ==================================================================
# 작업 자동 저장 / 복원
# ==================================================================
def _ensure_dirs():
    os.makedirs(CACHE_DIR, exist_ok=True)


def _cache_path(key):
    return os.path.join(CACHE_DIR, f"{key[0]}__{key[1]}.png")


def save_project(state, store):
    """state: dict(excel_path, name_pt, use_title, fields_*, season_map, slides)"""
    _ensure_dirs()
    # 이미지 캐시 저장
    for key, it in store.items():
        if it.img_bytes:
            p = _cache_path(key)
            if not os.path.exists(p):
                try:
                    Image.open(io.BytesIO(it.img_bytes)).convert("RGB").save(p, "PNG")
                except Exception:
                    pass
    meta = {key: {"name": it.name, "origin": it.origin, "season": it.season,
                  "size": it.size, "price": it.price}
            for key, it in store.items()}
    data = {
        "excel_path": state.get("excel_path"),
        "name_pt": state.get("name_pt", 14),
        "use_title": state.get("use_title", False),
        "fields_order": state.get("fields_order", DEFAULT_FIELDS),
        "fields_on": state.get("fields_on", {k: True for k in DEFAULT_FIELDS}),
        "season_map": state.get("season_map", {}),
        "slide_w": state.get("slide_w", SLIDE_W),
        "slide_h": state.get("slide_h", SLIDE_H),
        "slides": [{"cols": s["cols"], "rows": s["rows"],
                    "items": [list(k) for k in s["items"]]}
                   for s in state.get("slides", [])],
        "store": {f"{k[0]}||{k[1]}": v for k, v in meta.items()},
    }
    with open(PROJECT, "w", encoding="utf-8") as f:
        json.dump(data, f, ensure_ascii=False, indent=1)


def load_project():
    if not os.path.exists(PROJECT):
        return None
    try:
        with open(PROJECT, "r", encoding="utf-8") as f:
            data = json.load(f)
    except Exception:
        return None
    store = {}
    for sk, meta in data.get("store", {}).items():
        sheet, row = sk.split("||")
        key = (sheet, int(row))
        img = None
        p = _cache_path(key)
        if os.path.exists(p):
            with open(p, "rb") as fp:
                img = fp.read()
        store[key] = Item(meta.get("name"), meta.get("origin"),
                          meta.get("season"), meta.get("size"),
                          meta.get("price"), img, sheet, int(row))
    slides = [{"cols": s["cols"], "rows": s["rows"],
               "items": [tuple(k) for k in s["items"]]}
              for s in data.get("slides", [])]
    order = data.get("fields_order") or list(DEFAULT_FIELDS)
    on = data.get("fields_on") or {k: True for k in DEFAULT_FIELDS}
    for k in DEFAULT_FIELDS:                 # 누락 키 보정
        order.append(k) if k not in order else None
        on.setdefault(k, True)
    state = {
        "excel_path": data.get("excel_path"),
        "name_pt": data.get("name_pt", 14),
        "use_title": data.get("use_title", False),
        "fields_order": order,
        "fields_on": on,
        "season_map": data.get("season_map", {}),
        "slide_w": data.get("slide_w", SLIDE_W),
        "slide_h": data.get("slide_h", SLIDE_H),
        "slides": slides,
    }
    return state, store


# ==================================================================
# GUI
# ==================================================================
def run_gui():
    import tkinter as tk
    from tkinter import ttk, filedialog, messagebox
    from PIL import ImageTk, ImageDraw

    THUMB = 64
    CB = 44          # 체크박스 크기 (기본 대비 약 4배)

    class App(tk.Tk):
        def __init__(self):
            super().__init__()
            self.title("엑셀 → PPT 품목 카탈로그 생성기  v3")
            self.geometry("1320x800")

            self.field_order = list(DEFAULT_FIELDS)        # 입력 항목 순서
            self.field_on = {k: True for k in DEFAULT_FIELDS}  # 항목 포함 여부
            self.season_map = {}                            # 공급기간 치환표
            self.slide_w = SLIDE_W                           # 슬라이드 가로(EMU)
            self.slide_h = SLIDE_H                           # 슬라이드 세로(EMU)
            self._focus_key = None                           # 마지막 클릭 품목(이미지 지정용)
            self.excel_path = None
            self.store = {}            # key -> Item (모든 시트 누적)
            self.lib_keys = []         # 라이브러리 표시 순서(현재 시트)
            self.thumbs = {}           # key -> PIL thumbnail
            self.row_imgs = {}         # iid -> PhotoImage (참조 유지)
            self.checked = set()       # 체크된 품목 key (복수선택)
            self.slides = []           # [{cols,rows,items:[key]}]
            self.cur_slide = None      # 선택된 슬라이드 index
            self._drag_anchor = None
            self._drag_val = True      # 드래그 시 체크/해제 방향
            self._prev_win = None      # 미리보기 창
            self._prev_idx = 0

            self._build()
            self._restore()

        # ---------------- 레이아웃 ----------------
        def _build(self):
            top = ttk.Frame(self, padding=6)
            top.pack(fill="x")
            ttk.Button(top, text="📂 엑셀 불러오기", command=self.open_excel).pack(side="left")
            ttk.Label(top, text="시트:").pack(side="left", padx=(10, 2))
            self.sheet_cb = ttk.Combobox(top, state="readonly", width=14)
            self.sheet_cb.pack(side="left")
            self.sheet_cb.bind("<<ComboboxSelected>>", lambda e: self.load_sheet())

            ttk.Label(top, text="글자pt:").pack(side="left", padx=(8, 2))
            self.name_pt = tk.IntVar(value=14)
            ttk.Spinbox(top, from_=8, to=28, width=4, textvariable=self.name_pt).pack(side="left")
            self.use_title = tk.BooleanVar(value=False)
            ttk.Checkbutton(top, text="슬라이드 제목(시트명)", variable=self.use_title,
                            command=self.autosave).pack(side="left", padx=8)

            ttk.Label(top, text="슬라이드cm 가로:").pack(side="left")
            self.sw_cm = tk.DoubleVar(value=round(self.slide_w / 360000, 1))
            ttk.Spinbox(top, from_=10, to=50, increment=0.5, width=5,
                        textvariable=self.sw_cm,
                        command=self.apply_slide_size).pack(side="left")
            ttk.Label(top, text="세로:").pack(side="left", padx=(4, 0))
            self.sh_cm = tk.DoubleVar(value=round(self.slide_h / 360000, 1))
            ttk.Spinbox(top, from_=10, to=50, increment=0.5, width=5,
                        textvariable=self.sh_cm,
                        command=self.apply_slide_size).pack(side="left")

            ttk.Button(top, text="💾 PPT 내보내기", command=self.export_ppt).pack(side="right")
            self.open_after = tk.BooleanVar(value=True)
            ttk.Checkbutton(top, text="내보낸 후 바로 열기",
                            variable=self.open_after).pack(side="right", padx=4)
            self.save_lbl = ttk.Label(top, text="", foreground="#2a7")
            self.save_lbl.pack(side="right", padx=8)

            body = ttk.Frame(self, padding=(6, 0))
            body.pack(fill="both", expand=True)

            # ----- 왼쪽: 품목 라이브러리 -----
            left = ttk.LabelFrame(body, text="품목 라이브러리 (클릭=체크, 드래그=여러 개 체크)", padding=6)
            left.pack(side="left", fill="both", expand=True)

            lb = ttk.Frame(left); lb.pack(fill="x")
            ttk.Button(lb, text="전체체크", command=self.sel_all).pack(side="left")
            ttk.Button(lb, text="전체해제", command=self.sel_none).pack(side="left", padx=4)
            ttk.Button(lb, text="🖼 이미지 지정", command=self.assign_image).pack(side="left")
            self.chk_lbl = ttk.Label(lb, text="체크: 0개", foreground="#2a7")
            self.chk_lbl.pack(side="right")

            self.tree = ttk.Treeview(left, columns=("name", "origin", "used"),
                                     show="tree headings", selectmode="none")
            self.tree.heading("#0", text="✔  이미지")
            self.tree.heading("name", text="품목명")
            self.tree.heading("origin", text="원산지")
            self.tree.heading("used", text="배치")
            self.tree.column("#0", width=130, anchor="w")
            self.tree.column("name", width=260)
            self.tree.column("origin", width=80, anchor="center")
            self.tree.column("used", width=90, anchor="center")
            self.tree.pack(side="left", fill="both", expand=True)
            sb = ttk.Scrollbar(left, orient="vertical", command=self.tree.yview)
            sb.pack(side="right", fill="y")
            self.tree.configure(yscrollcommand=sb.set)
            ttk.Style(self).configure("Treeview", rowheight=THUMB + 8)
            self.tree.bind("<ButtonPress-1>", self.on_press)
            self.tree.bind("<B1-Motion>", self.on_drag)

            # ----- 가운데: 추가 버튼 + 항목/치환 설정 -----
            mid = ttk.Frame(body, padding=6); mid.pack(side="left", fill="y")
            ttk.Button(mid, text="체크한 품목 ▶\n현재 슬라이드에 추가",
                       command=self.add_to_slide).pack(pady=(6, 4), ipady=6, fill="x")
            ttk.Button(mid, text="◀ 슬라이드에서\n선택 품목 제거",
                       command=self.remove_from_slide).pack(pady=4, ipady=6, fill="x")
            ttk.Button(mid, text="🔍 미리보기",
                       command=self.preview).pack(pady=4, ipady=4, fill="x")

            # 입력 항목 / 순서
            ff = ttk.LabelFrame(mid, text="입력 항목 (더블클릭=ON/OFF)", padding=5)
            ff.pack(fill="x", pady=(10, 4))
            self.fields_lb = tk.Listbox(ff, height=5, width=24, exportselection=False,
                                        activestyle="none")
            self.fields_lb.pack(side="left", fill="x", expand=True)
            self.fields_lb.bind("<Double-Button-1>", lambda e: self.toggle_field())
            fbtn = ttk.Frame(ff); fbtn.pack(side="right", fill="y")
            ttk.Button(fbtn, text="▲", width=3,
                       command=lambda: self.move_field(-1)).pack(pady=1)
            ttk.Button(fbtn, text="▼", width=3,
                       command=lambda: self.move_field(1)).pack(pady=1)
            ttk.Button(fbtn, text="ON/OFF", width=7,
                       command=self.toggle_field).pack(pady=(6, 1))

            # 공급가능기간 치환
            sf = ttk.LabelFrame(mid, text="공급기간 치환 (원본→표시)", padding=5)
            sf.pack(fill="x", pady=4)
            row1 = ttk.Frame(sf); row1.pack(fill="x")
            ttk.Label(row1, text="원본:").pack(side="left")
            self.smap_orig = ttk.Combobox(row1, width=16)
            self.smap_orig.pack(side="left", fill="x", expand=True)
            row2 = ttk.Frame(sf); row2.pack(fill="x", pady=2)
            ttk.Label(row2, text="표시:").pack(side="left")
            self.smap_repl = ttk.Entry(row2, width=16)
            self.smap_repl.pack(side="left", fill="x", expand=True)
            ttk.Button(sf, text="추가 / 수정",
                       command=self.add_season_map).pack(fill="x", pady=2)
            self.smap_lb = tk.Listbox(sf, height=4, exportselection=False)
            self.smap_lb.pack(fill="x")
            ttk.Button(sf, text="선택 삭제",
                       command=self.del_season_map).pack(fill="x", pady=(2, 0))

            # ----- 오른쪽: 슬라이드 구성 -----
            right = ttk.LabelFrame(body, text="슬라이드 구성", padding=6)
            right.pack(side="right", fill="both", expand=True)

            sbar = ttk.Frame(right); sbar.pack(fill="x")
            ttk.Button(sbar, text="＋ 슬라이드 추가", command=self.add_slide).pack(side="left")
            ttk.Button(sbar, text="− 삭제", command=self.del_slide).pack(side="left", padx=4)
            ttk.Button(sbar, text="▲", width=3, command=lambda: self.move_slide(-1)).pack(side="left")
            ttk.Button(sbar, text="▼", width=3, command=lambda: self.move_slide(1)).pack(side="left")

            self.slide_lb = tk.Listbox(right, height=7, exportselection=False)
            self.slide_lb.pack(fill="x", pady=4)
            self.slide_lb.bind("<<ListboxSelect>>", lambda e: self.on_slide_select())

            grid_f = ttk.Frame(right); grid_f.pack(fill="x", pady=2)
            ttk.Label(grid_f, text="가로(열):").pack(side="left")
            self.cols_var = tk.IntVar(value=3)
            ttk.Spinbox(grid_f, from_=1, to=6, width=4, textvariable=self.cols_var,
                        command=self.apply_grid).pack(side="left", padx=(2, 10))
            ttk.Label(grid_f, text="세로(행):").pack(side="left")
            self.rows_var = tk.IntVar(value=2)
            ttk.Spinbox(grid_f, from_=1, to=6, width=4, textvariable=self.rows_var,
                        command=self.apply_grid).pack(side="left", padx=2)
            self.cap_lbl = ttk.Label(grid_f, text="", foreground="#888")
            self.cap_lbl.pack(side="left", padx=10)

            ttk.Label(right, text="이 슬라이드에 들어간 품목 (순서대로):").pack(anchor="w", pady=(6, 0))
            self.slide_items = tk.Listbox(right, height=10, exportselection=False,
                                          selectmode="extended")
            self.slide_items.pack(fill="both", expand=True)
            ib = ttk.Frame(right); ib.pack(fill="x")
            ttk.Button(ib, text="▲ 위로", command=lambda: self.move_item(-1)).pack(side="left")
            ttk.Button(ib, text="▼ 아래로", command=lambda: self.move_item(1)).pack(side="left", padx=4)

            self.status = tk.StringVar(value="엑셀을 불러오거나 이전 작업이 자동 복원됩니다.")
            ttk.Label(self, textvariable=self.status, relief="sunken", anchor="w",
                      padding=4).pack(fill="x", side="bottom")

            self.fill_fields()
            self.fill_season_map()

        # ---------------- 체크박스 이미지 ----------------
        def _row_image(self, key, checked):
            thumb = self.thumbs.get(key)
            tw = thumb.width if thumb else 60
            th = thumb.height if thumb else 60
            H = max(CB, th)
            W = CB + 8 + tw
            img = Image.new("RGBA", (W, H), (255, 255, 255, 0))
            d = ImageDraw.Draw(img)
            y0 = (H - CB) // 2
            if checked:
                d.rounded_rectangle([2, y0, CB - 2, y0 + CB - 4], radius=7,
                                    fill=(46, 134, 222), outline=(46, 134, 222), width=3)
                d.line([(CB * 0.24, y0 + CB * 0.52), (CB * 0.42, y0 + CB * 0.70)],
                       fill="white", width=5)
                d.line([(CB * 0.42, y0 + CB * 0.70), (CB * 0.76, y0 + CB * 0.28)],
                       fill="white", width=5)
            else:
                d.rounded_rectangle([2, y0, CB - 2, y0 + CB - 4], radius=7,
                                    outline=(120, 120, 120), width=3, fill=(255, 255, 255))
            if thumb:
                img.paste(thumb, (CB + 8, (H - th) // 2))
            return ImageTk.PhotoImage(img)

        def update_row(self, iid):
            key = self._iid_key(iid)
            photo = self._row_image(key, key in self.checked)
            self.row_imgs[iid] = photo
            self.tree.item(iid, image=photo)

        def update_all_rows(self):
            for iid in self.tree.get_children():
                self.update_row(iid)
            self.chk_lbl.config(text=f"체크: {len(self.checked)}개")

        # ---------------- 키 ↔ iid ----------------
        def _iid_key(self, iid):
            sheet, row = iid.rsplit("##", 1)
            return (sheet, int(row))

        def _key_iid(self, key):
            return f"{key[0]}##{key[1]}"

        # ---------------- 입력 항목 (체크/순서) ----------------
        def fill_fields(self):
            self.fields_lb.delete(0, "end")
            for k in self.field_order:
                mark = "☑" if self.field_on.get(k, True) else "☐"
                self.fields_lb.insert("end", f"{mark} {FIELD_LABELS.get(k, k)}")

        def toggle_field(self):
            sel = self.fields_lb.curselection()
            if not sel:
                return
            k = self.field_order[sel[0]]
            self.field_on[k] = not self.field_on.get(k, True)
            self.fill_fields()
            self.fields_lb.selection_set(sel[0])
            self.autosave()
            self._sync_preview()

        def move_field(self, delta):
            sel = self.fields_lb.curselection()
            if not sel:
                return
            i = sel[0]; j = i + delta
            if not (0 <= j < len(self.field_order)):
                return
            self.field_order[i], self.field_order[j] = \
                self.field_order[j], self.field_order[i]
            self.fill_fields()
            self.fields_lb.selection_set(j)
            self.autosave()
            self._sync_preview()

        def _enabled_fields(self):
            return [k for k in self.field_order if self.field_on.get(k, True)]

        # ---------------- 공급기간 치환 ----------------
        def refresh_season_choices(self):
            vals = sorted({it.season for it in self.store.values() if it.season})
            self.smap_orig["values"] = vals

        def fill_season_map(self):
            self.smap_lb.delete(0, "end")
            for orig, repl in self.season_map.items():
                self.smap_lb.insert("end", f"{orig}  →  {repl}")

        def add_season_map(self):
            orig = self.smap_orig.get().strip()
            repl = self.smap_repl.get().strip()
            if not orig:
                return
            self.season_map[orig] = repl
            self.fill_season_map()
            self.smap_repl.delete(0, "end")
            self.autosave()
            self._sync_preview()

        def del_season_map(self):
            sel = self.smap_lb.curselection()
            if not sel:
                return
            orig = list(self.season_map.keys())[sel[0]]
            del self.season_map[orig]
            self.fill_season_map()
            self.autosave()
            self._sync_preview()

        # ---------------- 엑셀 ----------------
        def open_excel(self):
            path = filedialog.askopenfilename(
                title="엑셀 선택", filetypes=[("Excel", "*.xlsx"), ("All", "*.*")])
            if not path:
                return
            try:
                sheets = list_sheets(path)
            except Exception as e:
                messagebox.showerror("오류", f"엑셀을 열 수 없습니다:\n{e}")
                return
            self.excel_path = path
            self.sheet_cb["values"] = sheets
            if sheets:
                self.sheet_cb.current(0)
                self.load_sheet()

        def load_sheet(self):
            if not self.excel_path:
                return
            sheet = self.sheet_cb.get()
            self.status.set(f"'{sheet}' 불러오는 중...")
            self.update_idletasks()
            try:
                items = load_items(self.excel_path, sheet)
            except Exception as e:
                messagebox.showerror("오류", f"시트 로드 실패:\n{e}")
                return
            self.lib_keys = []
            for it in items:
                self.store[it.key] = it
                self.lib_keys.append(it.key)
            self._make_thumbs(items)
            self.fill_library()
            self.refresh_season_choices()
            self.autosave()
            withimg = sum(1 for it in items if it.img_bytes)
            self.status.set(f"'{sheet}' — {len(items)}개 품목, 이미지 {withimg}개 인식")
            if items and withimg == 0:
                messagebox.showwarning(
                    "이미지 인식 안 됨",
                    "이 시트에서 셀에 박힌 이미지를 찾지 못했습니다.\n\n"
                    "· 그림이 '셀 안에 삽입(IMAGE 함수/셀 위 떠있는 그림)'인지 확인\n"
                    "· .xls가 아니라 .xlsx 인지 확인\n"
                    "· 그림이 도형 그룹/연결(LINK)된 경우 인식되지 않을 수 있습니다.\n"
                    "텍스트(품목명/가격 등)는 정상 입력됩니다.")

        def _make_thumbs(self, items):
            for it in items:
                if not it.img_bytes or it.key in self.thumbs:
                    continue
                try:
                    im = Image.open(io.BytesIO(it.img_bytes)).convert("RGB")
                    im.thumbnail((THUMB, THUMB))
                    self.thumbs[it.key] = im
                except Exception:
                    pass

        # ---------------- 라이브러리 그리기 ----------------
        def fill_library(self):
            self.tree.delete(*self.tree.get_children())
            self.row_imgs.clear()
            for key in self.lib_keys:
                it = self.store.get(key)
                if not it:
                    continue
                iid = self._key_iid(key)
                photo = self._row_image(key, key in self.checked)
                self.row_imgs[iid] = photo
                placed = self._slides_with(key)
                used = ",".join(str(i + 1) for i in placed) if placed else ""
                self.tree.insert("", "end", iid=iid, image=photo,
                                 values=(it.name, it.origin, used))
            self.chk_lbl.config(text=f"체크: {len(self.checked)}개")

        def _slides_with(self, key):
            return [i for i, s in enumerate(self.slides) if key in s["items"]]

        # ---------------- 슬라이드 크기 ----------------
        def apply_slide_size(self):
            try:
                self.slide_w = int(round(float(self.sw_cm.get()) * 360000))
                self.slide_h = int(round(float(self.sh_cm.get()) * 360000))
            except Exception:
                return
            self.autosave()
            self._sync_preview()

        # ---------------- 이미지 직접 지정 ----------------
        def assign_image(self):
            key = self._focus_key
            if not key or key not in self.store:
                messagebox.showinfo("알림",
                                    "먼저 왼쪽 목록에서 이미지를 넣을 품목을 한 번 클릭하세요.")
                return
            it = self.store[key]
            path = filedialog.askopenfilename(
                title=f"'{it.name}' 이미지 선택",
                filetypes=[("이미지", "*.png *.jpg *.jpeg *.gif *.bmp"), ("All", "*.*")])
            if not path:
                return
            try:
                with open(path, "rb") as f:
                    data = f.read()
                im = Image.open(io.BytesIO(data)).convert("RGB")
                buf = io.BytesIO()
                im.save(buf, "PNG")
                it.img_bytes = buf.getvalue()
                thumb = im.copy()
                thumb.thumbnail((THUMB, THUMB))
                self.thumbs[key] = thumb
                _ensure_dirs()                       # 캐시에 즉시 저장(복원용)
                im.save(_cache_path(key), "PNG")
            except Exception as e:
                messagebox.showerror("오류", f"이미지를 읽을 수 없습니다:\n{e}")
                return
            self.fill_library()
            self.autosave()
            self._sync_preview()
            self.status.set(f"'{it.name}'에 이미지 지정 완료")

        # ---------------- 체크 토글 / 드래그 ----------------
        def on_press(self, event):
            iid = self.tree.identify_row(event.y)
            if not iid:
                return
            self._drag_anchor = iid
            key = self._iid_key(iid)
            self._focus_key = key             # 이미지 지정 대상(마지막 클릭)
            if key in self.checked:           # 토글
                self.checked.discard(key)
                self._drag_val = False
            else:
                self.checked.add(key)
                self._drag_val = True
            self.update_row(iid)
            self.chk_lbl.config(text=f"체크: {len(self.checked)}개")

        def on_drag(self, event):
            iid = self.tree.identify_row(event.y)
            if not iid or not self._drag_anchor:
                return
            children = list(self.tree.get_children())
            try:
                a = children.index(self._drag_anchor)
                b = children.index(iid)
            except ValueError:
                return
            lo, hi = sorted((a, b))
            for c in children[lo:hi + 1]:
                k = self._iid_key(c)
                if self._drag_val:
                    self.checked.add(k)
                else:
                    self.checked.discard(k)
                self.update_row(c)
            self.chk_lbl.config(text=f"체크: {len(self.checked)}개")

        def sel_all(self):
            self.checked = set(self.lib_keys)
            self.update_all_rows()

        def sel_none(self):
            self.checked.clear()
            self.update_all_rows()

        def _selected_keys(self):
            # 라이브러리 표시 순서대로 체크된 품목 반환
            return [k for k in self.lib_keys if k in self.checked]

        # ---------------- 슬라이드 ----------------
        def add_slide(self):
            self.slides.append({"cols": self.cols_var.get(),
                                "rows": self.rows_var.get(), "items": []})
            self.refresh_slides()
            self.slide_lb.selection_clear(0, "end")
            self.slide_lb.selection_set("end")
            self.on_slide_select()
            self.autosave()

        def del_slide(self):
            if self.cur_slide is None:
                return
            del self.slides[self.cur_slide]
            self.cur_slide = None
            self.refresh_slides()
            self.fill_library()
            self.autosave()
            self._sync_preview()

        def move_slide(self, delta):
            i = self.cur_slide
            if i is None:
                return
            j = i + delta
            if not (0 <= j < len(self.slides)):
                return
            self.slides[i], self.slides[j] = self.slides[j], self.slides[i]
            self.cur_slide = j
            self.refresh_slides()
            self.slide_lb.selection_set(j)
            self.on_slide_select()
            self.autosave()
            self._sync_preview(j)

        def refresh_slides(self):
            self.slide_lb.delete(0, "end")
            for i, s in enumerate(self.slides):
                cap = s["cols"] * s["rows"]
                self.slide_lb.insert(
                    "end", f"슬라이드 {i+1}  ({s['cols']}×{s['rows']})  "
                           f"{len(s['items'])}/{cap}")

        def on_slide_select(self):
            sel = self.slide_lb.curselection()
            if not sel:
                self.cur_slide = None
                return
            self.cur_slide = sel[0]
            s = self.slides[self.cur_slide]
            self.cols_var.set(s["cols"])
            self.rows_var.set(s["rows"])
            self.fill_slide_items()

        def apply_grid(self):
            if self.cur_slide is None:
                return
            s = self.slides[self.cur_slide]
            s["cols"] = self.cols_var.get()
            s["rows"] = self.rows_var.get()
            self.refresh_slides()
            self.slide_lb.selection_set(self.cur_slide)
            self.fill_slide_items()
            self.autosave()
            self._sync_preview(self.cur_slide)

        def fill_slide_items(self):
            self.slide_items.delete(0, "end")
            if self.cur_slide is None:
                self.cap_lbl.config(text="")
                return
            s = self.slides[self.cur_slide]
            cap = s["cols"] * s["rows"]
            self.cap_lbl.config(text=f"최대 {cap}개")
            for key in s["items"]:
                it = self.store.get(key)
                self.slide_items.insert("end", it.name if it else str(key))

        def add_to_slide(self):
            keys = self._selected_keys()
            if not keys:
                messagebox.showinfo("알림", "먼저 왼쪽에서 품목을 체크(클릭/드래그)하세요.")
                return
            # 슬라이드가 없으면 자동 생성, 선택 안 됐으면 마지막 슬라이드 사용
            if not self.slides:
                self.add_slide()
            if self.cur_slide is None:
                self.cur_slide = len(self.slides) - 1
                self.slide_lb.selection_clear(0, "end")
                self.slide_lb.selection_set(self.cur_slide)
                self.on_slide_select()
            s = self.slides[self.cur_slide]
            cap = s["cols"] * s["rows"]
            added, full = 0, False
            for key in keys:
                if len(s["items"]) >= cap:
                    full = True
                    break
                if key not in s["items"]:
                    s["items"].append(key)
                    self.checked.discard(key)
                    added += 1
            self.refresh_slides()
            self.slide_lb.selection_set(self.cur_slide)
            self.fill_slide_items()
            self.fill_library()
            self.autosave()
            self.status.set(f"{added}개 품목을 슬라이드 {self.cur_slide+1}에 추가"
                            + (f" (남은 {len(keys)-added}개는 칸 부족)" if full else ""))
            self.preview()                       # 추가할 때마다 미리보기 자동 표시
            self._sync_preview(self.cur_slide)
            if full:
                messagebox.showinfo("가득 참",
                                    f"슬라이드 {self.cur_slide+1}은 최대 {cap}개입니다.\n"
                                    f"{added}개만 추가됨. 새 슬라이드를 추가해 나머지를 넣으세요.")

        def remove_from_slide(self):
            if self.cur_slide is None:
                return
            sel = list(self.slide_items.curselection())
            if not sel:
                return
            s = self.slides[self.cur_slide]
            for idx in sorted(sel, reverse=True):
                del s["items"][idx]
            self.refresh_slides()
            self.slide_lb.selection_set(self.cur_slide)
            self.fill_slide_items()
            self.fill_library()
            self.autosave()
            self._sync_preview(self.cur_slide)

        def move_item(self, delta):
            if self.cur_slide is None:
                return
            sel = self.slide_items.curselection()
            if not sel:
                return
            i = sel[0]; j = i + delta
            s = self.slides[self.cur_slide]
            if not (0 <= j < len(s["items"])):
                return
            s["items"][i], s["items"][j] = s["items"][j], s["items"][i]
            self.fill_slide_items()
            self.slide_items.selection_set(j)
            self.autosave()
            self._sync_preview(self.cur_slide)

        # ---------------- 자동 저장 / 복원 ----------------
        def _state(self):
            return {
                "excel_path": self.excel_path,
                "name_pt": self.name_pt.get(),
                "use_title": self.use_title.get(),
                "fields_order": self.field_order,
                "fields_on": self.field_on,
                "season_map": self.season_map,
                "slide_w": self.slide_w,
                "slide_h": self.slide_h,
                "slides": self.slides,
            }

        def autosave(self):
            try:
                save_project(self._state(), self.store)
                self.save_lbl.config(text="● 자동저장됨")
                self.after(1500, lambda: self.save_lbl.config(text=""))
            except Exception:
                pass

        def _restore(self):
            res = load_project()
            if not res:
                return
            state, store = res
            self.store = store
            self.excel_path = state.get("excel_path")
            self.name_pt.set(state.get("name_pt", 14))
            self.use_title.set(state.get("use_title", False))
            self.field_order = state.get("fields_order", list(DEFAULT_FIELDS))
            self.field_on = state.get("fields_on", {k: True for k in DEFAULT_FIELDS})
            self.season_map = state.get("season_map", {})
            self.slide_w = state.get("slide_w", SLIDE_W)
            self.slide_h = state.get("slide_h", SLIDE_H)
            self.sw_cm.set(round(self.slide_w / 360000, 1))
            self.sh_cm.set(round(self.slide_h / 360000, 1))
            self.fill_fields()
            self.fill_season_map()
            self.slides = state.get("slides", [])
            # 썸네일 복원
            for key, it in store.items():
                if it.img_bytes and key not in self.thumbs:
                    try:
                        im = Image.open(io.BytesIO(it.img_bytes)).convert("RGB")
                        im.thumbnail((THUMB, THUMB))
                        self.thumbs[key] = im
                    except Exception:
                        pass
            # 라이브러리: 엑셀 있으면 시트 목록 복원
            if self.excel_path and os.path.exists(self.excel_path):
                try:
                    self.sheet_cb["values"] = list_sheets(self.excel_path)
                    if self.sheet_cb["values"]:
                        self.sheet_cb.current(0)
                        self.load_sheet()
                except Exception:
                    pass
            else:
                # 엑셀 없으면 캐시에 있는 품목 전부 라이브러리에 표시
                self.lib_keys = list(store.keys())
                self.fill_library()
            self.refresh_season_choices()
            self.refresh_slides()
            self.status.set("이전 작업을 복원했습니다.")

        # ---------------- 미리보기 ----------------
        def _titles(self):
            if self.use_title.get():
                return [self.sheet_cb.get() or ""] * len(self.slides)
            return None

        def preview(self):
            if not self.slides:
                messagebox.showinfo("알림", "먼저 슬라이드를 추가하세요.")
                return
            self._prev_idx = self.cur_slide if self.cur_slide is not None else 0
            if self._prev_win and tk.Toplevel.winfo_exists(self._prev_win):
                self._prev_win.lift()
            else:
                self._prev_win = tk.Toplevel(self)
                self._prev_win.title("미리보기")
                nav = ttk.Frame(self._prev_win); nav.pack(fill="x", pady=4)
                ttk.Button(nav, text="◀ 이전",
                           command=lambda: self._prev_nav(-1)).pack(side="left", padx=6)
                self._prev_lbl = ttk.Label(nav, text="")
                self._prev_lbl.pack(side="left", expand=True)
                ttk.Button(nav, text="다음 ▶",
                           command=lambda: self._prev_nav(1)).pack(side="right", padx=6)
                self._prev_canvas = ttk.Label(self._prev_win)
                self._prev_canvas.pack(padx=6, pady=6)
            self._prev_render()

        def _prev_nav(self, d):
            self._prev_idx = max(0, min(len(self.slides) - 1, self._prev_idx + d))
            self._prev_render()

        def _prev_open(self):
            return bool(self._prev_win) and tk.Toplevel.winfo_exists(self._prev_win)

        def _sync_preview(self, idx=None):
            """슬라이드 변경 시 미리보기 창이 열려 있으면 자동 갱신."""
            if not self._prev_open() or not self.slides:
                return
            if idx is not None:
                self._prev_idx = max(0, min(len(self.slides) - 1, idx))
            elif self._prev_idx >= len(self.slides):
                self._prev_idx = len(self.slides) - 1
            self._prev_render()

        def _prev_render(self):
            from PIL import ImageTk
            i = self._prev_idx
            sd = self.slides[i]
            titles = self._titles()
            title = titles[i] if titles else None
            img = render_slide_image(sd, self.store, title=title,
                                     fields=self._enabled_fields(),
                                     season_map=self.season_map, width=900,
                                     slide_w=self.slide_w, slide_h=self.slide_h)
            self._prev_photo = ImageTk.PhotoImage(img)
            self._prev_canvas.config(image=self._prev_photo)
            self._prev_lbl.config(
                text=f"슬라이드 {i+1} / {len(self.slides)}   "
                     f"({sd['cols']}×{sd['rows']}, {len(sd['items'])}개)")

        # ---------------- 내보내기 ----------------
        def export_ppt(self):
            if not self.slides or all(not s["items"] for s in self.slides):
                messagebox.showwarning("알림", "슬라이드에 품목을 먼저 추가하세요.")
                return
            path = filedialog.asksaveasfilename(
                title="PPT 저장", defaultextension=".pptx",
                filetypes=[("PowerPoint", "*.pptx")],
                initialfile="품목_카탈로그.pptx")
            if not path:
                return
            try:
                n = build_pptx(self.slides, self.store, path,
                               fields=self._enabled_fields(),
                               season_map=self.season_map,
                               name_pt=self.name_pt.get(), titles=self._titles(),
                               slide_w=self.slide_w, slide_h=self.slide_h)
            except Exception as e:
                messagebox.showerror("오류", f"PPT 생성 실패:\n{e}")
                return
            self.status.set(f"완료: {n}개 슬라이드 → {path}")
            if self.open_after.get():
                if not open_file(path):
                    messagebox.showinfo("완료", f"{n}개 슬라이드 생성\n{path}\n(자동 열기는 실패)")
            else:
                messagebox.showinfo("완료", f"{n}개 슬라이드 생성\n{path}")

    App().mainloop()


if __name__ == "__main__":
    run_gui()
